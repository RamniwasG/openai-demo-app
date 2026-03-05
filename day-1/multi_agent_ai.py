from agents import Agent, Runner
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
import datetime

load_dotenv()


# -----------------------------
# 1️⃣ Define Agents
# -----------------------------

researcher_agent = Agent(
    name="Researcher AI",
    instructions="""
You are a research specialist.
Collect detailed and accurate information about the given topic.
Provide structured bullet points with key facts.
""",
    model="gpt-4.1-mini"
)


analyst_agent = Agent(
    name="Analyst AI",
    instructions="""
You are a data analyst.
Analyze the provided research.
Identify trends, insights, risks, opportunities.
Summarize findings clearly in bullet points.
""",
    model="gpt-4.1-mini"
)


writer_agent = Agent(
    name="Writer AI",
    instructions="""
You are a professional report writer.
Using the research and analysis provided,
write a structured professional report with:

- Title
- Introduction
- Key Findings
- Analysis
- Conclusion
- Recommendations

Keep tone formal and clear.
""",
    model="gpt-4.1-mini"
)


# -----------------------------
# 2️⃣ Manager Workflow
# -----------------------------

def manager_workflow(topic: str):

    print("\n🔎 Step 1: Researching...\n")

    research_response = Runner.run_sync(
        starting_agent=researcher_agent,
        input=topic
    )

    research_data = research_response.final_output


    print("\n📊 Step 2: Analyzing...\n")

    analysis_response = Runner.run_sync(
        starting_agent=analyst_agent,
        input=research_data
    )

    analysis_data = analysis_response.final_output


    print("\n✍️ Step 3: Writing Report...\n")

    combined_input = f"""
    Topic: {topic}

    Research Data:
    {research_data}

    Analysis Data:
    {analysis_data}
    """

    final_report_response = Runner.run_sync(
        starting_agent=writer_agent,
        input=combined_input
    )

    final_report = final_report_response.final_output

    return final_report


# -----------------------------
# 3️⃣ Generate TXT File
# -----------------------------
def generate_txt(report_text, filename):
    with open(filename, "w", encoding="utf-8") as f:
        f.write(report_text)
    print(f"\n✅ TXT report saved as {filename}")


# -----------------------------
# 4 Generate PPT File
# -----------------------------
def generate_ppt(slide_text: str, filename="AI_Report.pptx"):

    presentation = Presentation()

    slides = slide_text.split("Slide")

    for slide_data in slides:
        if slide_data.strip() == "":
            continue

        lines = slide_data.strip().split("\n")
        title = lines[0].replace(":", "").strip()
        content_lines = lines[1:]

        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        content = slide.placeholders[1]

        content.text = ""

        for line in content_lines:
            if line.strip():
                p = content.text_frame.add_paragraph()
                p.text = line.strip("- ").strip()
                p.level = 1

    presentation.save(filename)
    print(f"\n✅ PPT saved as {filename}")


# -----------------------------
# 5 Generate DOC File
# -----------------------------
def generate_docx(report_text: str, filename="AI_Report.docx"):

    document = Document()
    document.add_heading("AI Generated Report", level=1)

    for line in report_text.split("\n"):
        document.add_paragraph(line)

    document.save(filename)
    print(f"\n✅ Report saved as {filename}")


# -----------------------------
# 6 Main Entry
# -----------------------------

def main():

    print("\n=== Multi-Agent AI Report Generator ===")
    print("1 → TXT")
    print("2 → DOCX")
    print("3 → PPTX")
    print("Type 'quit' to exit\n")

    topic = input("Enter report topic: ")

    if topic.lower() == "quit":
        print("Exiting...")
        return

    file_choice = input("Choose file type (1/2/3): ")

    final_report = manager_workflow(topic)

    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    base_filename = f"AI_Report_{timestamp}"

    if file_choice == "1":
        generate_txt(final_report, base_filename + ".txt")
    elif file_choice == "2":
        generate_docx(final_report, base_filename + ".docx")
    elif file_choice == "3":
        generate_ppt(final_report, base_filename + ".pptx")
    else:
        print("❌ Invalid option. No file generated.")

    print("\n📄 Report Preview:\n")
    print(final_report)



if __name__ == "__main__":
    main()